import argparse
from openai import OpenAI
import os
import json
import subprocess
import re
from docx import Document
from pptx import Presentation
from lxml import etree
import fitz
import base64
import hashlib
import hmac
import time
import requests
import urllib

lfasr_host = 'https://raasr.xfyun.cn/v2/api'
# 请求的接口名
api_upload = '/upload'
api_get_result = '/getResult'

## 设置API
api_key =  "阿里云百炼平台的API"
base_url = "https://dashscope.aliyuncs.com/compatible-mode/v1"
os.environ["OPENAI_API_KEY"]  =  api_key

client = OpenAI(
  api_key= api_key,
  base_url= base_url
)

prompt_dir = os.path.dirname(os.path.abspath(__file__))

def get_response(prompt):
  response = client.chat.completions.create(
    messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
    ],
    model="qwen-plus",
)
  return response.choices[0].message.content
  
def extract_json_from_string(input_string):
    # 使用正则表达式匹配JSON部分
    json_pattern = r"```json\n(.*?)```"
    match = re.search(json_pattern, input_string, re.DOTALL)

    if match:
        json_str = match.group(1).strip()
        # 转换为JSON对象
        json_data = json.loads(json_str)
        return json_data
    else:
        return None
#---------------- 可能需要修改的部分------------------------------
def extract_text_docx(file_path):
    doc = Document(file_path)
    full_text = []
    
    # 定义命名空间映射
    nsmap = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

    # 1. 提取普通段落文字
    for para in doc.paragraphs:
        full_text.append(para.text)
    
    # 2. 提取表格中的文字
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                full_text.append(cell.text)
    
    root = etree.fromstring(doc.element.xml)
    txbx_contents = root.xpath('.//w:txbxContent', namespaces=nsmap)
    for txbx in txbx_contents:
        # 查找文本框内的所有段落，并直接提取其中所有文本
        texts = txbx.xpath('.//w:t/text()', namespaces=nsmap)
        if texts:
            full_text.append("".join(texts))
    
    return "\n".join(full_text)

def extract_info_pdf(pdf_path, output_folder="extracted_images"):
    """
    提取 PDF 文档中的所有文字和图片。
    
    参数：
    - pdf_path: PDF 文件路径。
    - output_folder: 存储提取图片的文件夹，默认为 "extracted_images"。
    
    返回：
    - 所有页面的文本拼接成的字符串。
    """
    # 打开 PDF 文件
    doc = fitz.open(pdf_path)
    full_text = []

    # 确保图片存储文件夹存在
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    image_counter = 1

    # 遍历所有页面
    for page_index in range(len(doc)):
        page = doc.load_page(page_index)
        # 提取页面文本
        text = page.get_text()
        full_text.append(text)
        
        # 提取页面内的图片（使用 get_images(True) 可以获得更多信息）
        image_list = page.get_images(full=True)
        for img in image_list:
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = os.path.join(output_folder, f"image_page{page_index+1}_{image_counter}.{image_ext}")
            with open(image_filename, "wb") as img_file:
                img_file.write(image_bytes)
            image_counter += 1

    return "\n".join(full_text)

def extract_text_pptx(file_path):
    """从 PPT 文件 (.pptx) 中提取文本"""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)

def extract_text_from_file(file_path):
    """根据文件后缀自动选择相应的提取方法"""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.docx':
        return extract_text_docx(file_path)
    elif ext == '.pptx':
        return extract_text_pptx(file_path)
    elif ext == '.pdf':
        return extract_info_pdf(file_path)
    else:
        raise ValueError("仅支持 .docx , .pdf 和 .pptx 格式文件")

def save_text_to_file(text, output_path):
    """将提取的文本保存到 txt 文件中"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)
    print(f"文本已保存到 {output_path}")

def read_text_file(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

# 使用 OpenAI 模型提取知识点和关联关系
def extract_knowledge(path, text, style="tree"):
    """
    将文本传递给大语言模型，提取知识点和关联关系，要求返回格式如下：
    {
      "nodes": [
          {"id": "知识点名称", "type": "概念或其他描述"},
          ...
      ],
      "edges": [
          {"source": "知识点1", "target": "知识点2", "relation": "关联类型"},
          ...
      ]
    }
    """
    if style=="tree":
        sample = {
                    "name": "序列决策问题",
                    "child": [
                              {"name": "状态", "子节点": []},{
                                "节点": "智能体与环境交互",
                                    "子节点": [
                                    {"节点": "状态", "子节点": []},
                                    {"节点": "动作", "子节点": []},
                                    {"节点": "奖励", "子节点": []}]
                        },
                        {
                        "节点": "决策目标",
                                "子节点": [
                                        {"节点": "最大化累计回报", "子节点": []}
                                          ]
                    }
                    ]
                    }
        attributes = ["重点", "难点", "案例", "总结", "考点", "概述", "实操/训练", "练习", "问题(引例）", "项目/任务/步骤", "外延", "讨论", "情景引入",
    "实验", "岗位", "证书", "比赛"]
        prompt = (f"""请根据以下文本内容，构建一棵树状的知识图谱。要求：
                    1. 顶层节点为本章的主要主题（例如“序列决策问题”、“马尔可夫决策过程”、“值迭代”等）。
                    2. 每个顶层节点下分为若干子节点，表示该主题下的关键概念、定义、模型、方法或例子。
                    3. 每个子节都包含一个属性，属性包括了{attributes}等。
                    4. 子节点下可以继续分层，展示更细粒度的知识点或解释。
                    5. 只需要提取所给文字内容包含的知识点和解释，不必添加额外的知识点和解释。
                    6. 输出格式采用 JSON 树结构，形如：
                    {sample}
                    
                    请根据下面的文本内容生成知识图谱：{text}""")

    else:
        sample = {
                        "nodes": [
                        {"id": "1", "name": "序列决策问题"},
                        {"id": "2", "name": "状态"},
                        {"id": "3", "name": "动作"},
                        {"id": "4", "name": "奖励"}
                     ],
                        "edges": [
                        {"source": "1", "relation": "包含", "target": "状态"},
                        {"source": "1", "relation": "包含", "target": "动作"},
                        {"source": "1", "relation": "包含", "target": "奖励"},
                        {"source": "状态", "relation": "影响", "target": "奖励"}
                        ]
                        }
        prompt = (f"""
                    请根据以下文本内容抽取所有核心知识点，并建立它们之间的连接关系，构造一个网状的知识图谱。要求：
                    1. 从文本中识别出所有重要概念、定义、模型、算法、例子等作为节点。
                    2. 针对每对知识点，识别它们之间的关系，如“属于”、“包含”、“导致”、“相互影响”等。
                    3. 输出格式为 JSON 格式，包含两个部分：一个是所有节点列表（每个节点含有唯一ID和名称），另一个是所有关系列表（每个关系包含“起始节点ID”、“关系类型”、“结束节点ID”）。
                    例如：
                        {sample}   

                        请根据下面的文本内容生成知识图谱：{text}""")
    
    response = client.chat.completions.create(
        model="qwen-max",  # 或其他你有权限使用的模型
        messages=[
            {"role": "system", "content": "你是一个知识图谱构建专家。"},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
        #max_tokens=1024
    )
    #print(response)
    result = response.choices[0].message.content.strip()
    with open(os.path.join(path, 'tree1.json'), 'w', encoding = "utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=4)
# #---------------- 可能需要修改的部分------------------------------

def generate_thumbnail(input_path, output_path):
    """生成视频缩略图"""
    try:
        subprocess.run([
            'ffmpeg',
            '-i', input_path,
            '-ss', '00:00:01',
            '-vframes', '1',
            '-q:v', '2',
            output_path
        ], check=True)
    except Exception as e:
        print(f"生成缩略图失败: {str(e)}")

def generate_audio(path):
    """
    批量将视频文件转换为音频文件。
    
    参数:
        video_dir (str): 包含视频文件的目录。
        output_dir (str): 输出音频文件的目录。
    """
    # 检查输出目录是否存在，如果不存在则创建
    if not os.path.exists(path):
        os.makedirs(path)

    for filename in os.listdir(path):
        if filename.endswith(('.mp4', '.avi', '.mov', '.mkv')):  # 支持的文件格式
            video_path = os.path.join(path, filename)
            audio_path = os.path.join(path, 'audio.mp3')

            # 构建 FFmpeg 命令
            command = [
                'ffmpeg',
                '-i', video_path,
                '-vn',  # 不包括视频
                '-acodec', 'libmp3lame',  # 使用 MP3 编解码器
                '-ab', '128k',  # 设置音频比特率
                audio_path  # 输出文件路径
            ]

            # 执行 FFmpeg 命令
            try:
                subprocess.run(command, check=True)
                print(f"转换完成：{audio_path}")
            except subprocess.CalledProcessError as e:
                print(f"转换失败：{filename}，错误：{e}")
    return

def convert_to_srt(data):
    srt_lines = []
    index = 1

    for lattice in data['lattice2']:
        begin = int(lattice['begin'])
        end = int(lattice['end'])
        sentence = []
        start_time = begin
        end_time = end

        for segment in lattice['json_1best']['st']['rt']:
            for word_segment in segment['ws']:
                sentence.append(''.join([word['w'] for word in word_segment['cw']]))

        # Join the sentence and format the time
        full_sentence = ''.join(sentence)
        start_time_srt = f"{start_time // 3600000:02d}:{(start_time % 3600000) // 60000:02d}:{(start_time % 60000) // 1000:02d},{start_time % 1000:03d}"
        end_time_srt = f"{end_time // 3600000:02d}:{(end_time % 3600000) // 60000:02d}:{(end_time % 60000) // 1000:02d},{end_time % 1000:03d}"

        srt_lines.append(f"{index}\n--> speaker: \n{start_time_srt} --> {end_time_srt}\n{full_sentence}\n\n")
        index += 1

    return ''.join(srt_lines)

class RequestApi(object):
    def __init__(self, appid, secret_key, upload_file_path):
        self.appid = appid
        self.secret_key = secret_key
        self.upload_file_path = upload_file_path
        self.ts = str(int(time.time()))
        self.signa = self.get_signa()

    def get_signa(self):
        appid = self.appid
        secret_key = self.secret_key
        m2 = hashlib.md5()
        m2.update((appid + self.ts).encode('utf-8'))
        md5 = m2.hexdigest()
        md5 = bytes(md5, encoding='utf-8')
        # 以secret_key为key, 上面的md5为msg， 使用hashlib.sha1加密结果为signa
        signa = hmac.new(secret_key.encode('utf-8'), md5, hashlib.sha1).digest()
        signa = base64.b64encode(signa)
        signa = str(signa, 'utf-8')
        return signa

    def upload(self):
        upload_file_path = self.upload_file_path
        file_len = os.path.getsize(upload_file_path)
        file_name = os.path.basename(upload_file_path)

        param_dict = {}
        param_dict['appId'] = self.appid
        param_dict['signa'] = self.signa
        param_dict['ts'] = self.ts
        param_dict["fileSize"] = file_len
        param_dict["fileName"] = file_name
        param_dict["duration"] = "200"
        data = open(upload_file_path, 'rb').read(file_len)

        response = requests.post(url=lfasr_host + api_upload + "?" + urllib.parse.urlencode(param_dict),
                                 headers={"Content-type": "application/json"}, data=data)
        result = json.loads(response.text)
        return result

    def get_result(self):
        uploadresp = self.upload()
        orderId = uploadresp['content']['orderId']
        param_dict = {}
        param_dict['appId'] = self.appid
        param_dict['signa'] = self.signa
        param_dict['ts'] = self.ts
        param_dict['orderId'] = orderId
        param_dict['resultType'] = "transfer,predict"
        status = 3
        # 建议使用回调的方式查询结果，查询接口有请求频率限制
        while status == 3:
            response = requests.post(url=lfasr_host + api_get_result + "?" + urllib.parse.urlencode(param_dict),
                                     headers={"Content-type": "application/json"})
            result = json.loads(response.text)
            status = result['content']['orderInfo']['status']
            if status == 4:
                break
            time.sleep(5)
        return result

def generate_video_tree(path):
    # 打开文件并按行读取内容
    with open(os.path.join(path, 'subtitles.srt'), 'r', encoding='utf-8') as file:
        lines = file.readlines()
    with open(os.path.join(prompt_dir, 'prompt1.txt'), 'r', encoding='utf-8') as file:
        sample = file.read()
    # 构建prompt
    prompt = [
        sample,
      "\n教学内容如下：\n"
    ] + lines
    prompt = "".join(prompt)
    response = get_response(prompt)
    response = extract_json_from_string(response)
    with open(os.path.join(path, 'tree1.json'), 'w', encoding = "utf-8") as f:
        json.dump(response, f, ensure_ascii=False, indent=4)
    return


def generate_document_tree(path):
    # 打开文件并按行读取内容
    supported_ext = ('.docx', '.pptx', '.pdf')
    all_text = []
    
    for filename in os.listdir(path):
        if filename.lower().endswith(supported_ext):
            file_path = os.path.join(path, filename)
            try:
                text = extract_text_from_file(file_path)
                all_text.append(f"=== 文件: {filename} ===\n{text}\n\n")
            except Exception as e:
                print(f"处理文件 {filename} 时出错: {str(e)}")
    
    output_file = "output.txt"
    if all_text:
        with open(os.path.join(path, output_file), 'w', encoding='utf-8') as f:
            f.write("\n".join(all_text))
        print(f"所有文本已提取并保存到 {output_file}")
    else:
        print("未找到支持的文档文件")
        with open(os.path.join(path, 'tree2.json'), 'w', encoding = "utf-8") as f:
            json.dump({}, f, ensure_ascii=False, indent=4)
            return
    file_path = os.path.join(path, 'outline.txt')
    text = read_text_file(file_path)

    response = extract_knowledge(text,"tree")
    with open(os.path.join(path, 'tree2.json'), 'w', encoding = "utf-8") as f:
        json.dump(response, f, ensure_ascii=False, indent=4)

def generate_report(path):
    # 打开文件并按行读取内容
    with open(os.path.join(path, 'tree1.json'), 'r', encoding='utf-8') as file: #修改成视频图谱所在的文件夹
        tree1 = file.readlines()
    with open(os.path.join(path, 'subtitles.srt'), 'r', encoding='utf-8') as file: #修改成字幕所在的文件夹
        srt = file.readlines()
    with open(os.path.join(prompt_dir, 'prompt3.txt'), 'r', encoding='utf-8') as file:
        prompt3 = file.read()
    with open(os.path.join(prompt_dir, 'prompt4.txt'), 'r', encoding='utf-8') as file:
        prompt4 = file.read()
    # 构建prompt
    prompt = f"""
              {prompt3}
              图谱如下：{tree1}
              字幕如下：{srt}
              """
    response1 = get_response(prompt)
    with open(os.path.join(path, 'tree2.json'), 'r', encoding='utf-8') as file: #修改成字幕所在的文件夹
        tree2 = file.readlines()
    prompt = (f"""{prompt4}
              这里是结构化知识信息:{tree2}
              这里是视频转录文字:{srt}""")
    
    response2 = get_response(prompt)
    with open(os.path.join(path, 'report.md'), 'w', encoding='utf-8') as f:
        f.write(response1)
        f.write(response2)
    return

def generate_outline(path):
    with open(os.path.join(prompt_dir, 'prompt5.txt'), 'r', encoding='utf-8') as file:
        prompt5 = file.read()
    with open(os.path.join(path, 'tree1.json'), 'r', encoding='utf-8') as file: #修改成视频图谱所在的文件夹
        tree1 = file.readlines()
    with open(os.path.join(path, 'subtitles.srt'), 'r', encoding='utf-8') as file: #修改成字幕所在的文件夹
        srt = file.readlines()
    # 构建prompt
    prompt = f"""
              {prompt5}
              图谱如下：{tree1}
              字幕如下：{srt}
              """
    response = get_response(prompt)
    with open(os.path.join(path, 'outline.md'), 'w', encoding='utf-8') as f:
        f.write(response)
    return


def process_video(input_path, output_dir, video_name):
    # 1. 移动视频文件
    # video_path = os.path.join(output_dir, 'video.mp4')
    # # os.rename(input_path, video_path)
    # # print('1')
    # # 2. 生成缩略图
    # generate_thumbnail(video_path, os.path.join(output_dir, 'thumbnail.jpg'))
    # print('封面图生成成功...')
    # # 3. 转录音频
    # generate_audio(output_dir)
    # print('音频转录成功...')
    # # 4. 转录字幕
    # api = RequestApi(appid="05813e94",
    #                  secret_key="424545dc683a117f387e2251a107e1bc",
    #                  upload_file_path=os.path.join(output_dir, 'audio.mp3'))
    # result = api.get_result()
    
    # # 解析嵌套的 JSON 字符串
    # order_result = json.loads(result['content']['orderResult'])
    # print("Parsed orderResult:", order_result)
    # order_result = convert_to_srt(order_result)

    # # 保存结果到文件
    # with open(os.path.join(output_dir, 'subtitles.srt'), 'w', encoding='utf-8') as f:
    #     f.write(order_result)
    # print('字幕转录成功...')
    # 5. 生成视频图谱
    generate_video_tree(output_dir)
    print('视频图谱生成成功...')
    # 6. 生成教案图谱

    generate_document_tree(output_dir)
    # with open(os.path.join(output_dir, 'tree2.json'), 'w', encoding='utf-8') as f:
    #     json.dump({}, f, indent=4)
    print('教案图谱生成成功...')
    # 7. 生成报告
    generate_report(output_dir)
    print('分析报告生成成功...')
    # 8. 生成新教案
    generate_outline(output_dir)
    print('新教案生成成功...')
    print(f"视频处理完成: {video_name}")

if __name__ == '__main__':
    # parser = argparse.ArgumentParser()
    # parser.add_argument('--input', required=True)
    # parser.add_argument('--output', required=True) # 根据用户输入名称新建的一个文件夹data/name
    # parser.add_argument('--name', required=True)
    # args = parser.parse_args()
    
    # os.makedirs(args.output, exist_ok=True)
    # process_video(args.input, args.output, args.name)
    process_video(None, 'backend/data/user1/English_translation1', 'video.mp4')
    process_video(None, 'backend/data/user1/English_translation2', 'video.mp4')
